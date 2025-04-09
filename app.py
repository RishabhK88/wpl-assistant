from omegaconf import OmegaConf
from query import VectaraQuery
import os
from PIL import Image
import uuid
import io
import pptx
from pptx.util import Inches
import pandas as pd
from fpdf import FPDF
import plotly.express as px
import plotly.io as pio
import time
from typing import List, Dict

import streamlit as st
from streamlit_pills import pills
from streamlit_feedback import streamlit_feedback

from utils import thumbs_feedback, send_amplitude_data, escape_dollars_outside_latex

from dotenv import load_dotenv
load_dotenv(override=True)

import debugpy

def initialize_debugger(max_attempts=30):

    try:
        if not hasattr(initialize_debugger, 'is_listening'):
            debugpy.listen(("0.0.0.0", 5678))
            initialize_debugger.is_listening = True
            print("🔌 Debugger listening on port 5678")

        if debugpy.is_client_connected():
            print("✅ Debugger already attached")
            return True

        print("⏳ Waiting for debugger to attach...")
        attempts = 0
        while not debugpy.is_client_connected() and attempts < max_attempts:
            time.sleep(1)
            attempts += 1

        if debugpy.is_client_connected():
            print("🎯 Debugger attached successfully!")
            return True
        else:
            print("⚠️ Debugger connection timeout")
            return False

    except Exception as e:
        return False

if os.getenv('DEBUGGER', 'False').lower() == 'true':
    initialize_debugger()

MAX_EXAMPLES = 6
LANGUAGES = {'English': 'eng', 'Spanish': 'spa', 'French': 'fra', 'Chinese': 'zho', 'German': 'deu', 'Hindi': 'hin', 'Arabic': 'ara',
             'Portuguese': 'por', 'Italian': 'ita', 'Japanese': 'jpn', 'Korean': 'kor', 'Russian': 'rus', 'Turkish': 'tur', 'Persian (Farsi)': 'fas',
             'Vietnamese': 'vie', 'Thai': 'tha', 'Hebrew': 'heb', 'Dutch': 'nld', 'Indonesian': 'ind', 'Polish': 'pol', 'Ukrainian': 'ukr',
             'Romanian': 'ron', 'Swedish': 'swe', 'Czech': 'ces', 'Greek': 'ell', 'Bengali': 'ben', 'Malay (or Malaysian)': 'msa', 'Urdu': 'urd'}
INITIAL_ASSISTANT_MESSAGE = "How may I help you?"
ASSISTANT_AVATAR = '🤖'
USER_AVATAR = '🧑‍💻'
LOGO_PATH = 'ti-logo.png'

def isTrue(x) -> bool:
    if isinstance(x, bool):
        return x
    return isinstance(x, str) and x.strip().lower() == 'true'

def export_to_ppt(content, fig=None):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    if isinstance(content, pd.DataFrame):
        title.text = "Data Table"
        try:
            rows, cols = content.shape
            top_inch = Inches(1.5) if fig is None else Inches(0.8)
            height_inch = Inches(5.0) if fig is None else Inches(2.5)
            table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), top_inch, Inches(9), height_inch).table
            for c_idx, col_name in enumerate(content.columns):
                table.cell(0, c_idx).text = str(col_name)
            for r_idx in range(rows):
                for c_idx in range(cols):
                    table.cell(r_idx + 1, c_idx).text = str(content.iloc[r_idx, c_idx])
        except Exception as e:
            body.text_frame.text = f"Could not automatically add table to PPT.\nError: {e}\n\nData:\n{content.to_string()}"

    else:
        title.text = "Response"
        body.text_frame.text = str(content)

    if fig is not None:
        if isinstance(content, pd.DataFrame):
            img_stream = io.BytesIO()
            pio.write_image(fig, img_stream, format='png', scale=2)
            img_stream.seek(0)
            pic_top = top_inch + height_inch + Inches(0.2)
            slide.shapes.add_picture(img_stream, Inches(0.5), pic_top, width=Inches(9))
        else:
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            title_shape.text_frame.text = "Visualization"
            img_stream = io.BytesIO()
            pio.write_image(fig, img_stream, format='png', scale=2)
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.0), width=Inches(9))

    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream.getvalue()

def export_to_pdf(content, fig=None):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    page_width = pdf.w - 2 * pdf.l_margin

    pdf.cell(page_width, 10, txt="Response:", ln=True, align='L')
    pdf.ln(5)

    if isinstance(content, pd.DataFrame):
        pdf.set_font("Arial", 'B', size=10)
        header = content.columns.tolist()
        col_widths = {}
        for col in header:
            col_widths[col] = pdf.get_string_width(str(col)) + 6

        for _, row in content.iterrows():
            for col in header:
                 width = pdf.get_string_width(str(row[col])) + 6
                 if width > col_widths[col]:
                     col_widths[col] = width

        total_width = sum(col_widths.values())
        scale_factor = 1.0
        if total_width > page_width:
             scale_factor = page_width / total_width

        final_widths = [col_widths[h] * scale_factor for h in header]

        for h, w in zip(header, final_widths):
            pdf.cell(w, 10, txt=str(h), border=1, align='C', fill=True)
        pdf.ln()

        pdf.set_font("Arial", size=10)
        for _, row in content.iterrows():
            for item, w in zip(row, final_widths):
                text = str(item)
                char_width_estimate = pdf.get_string_width('W') * 0.8
                if char_width_estimate > 0:
                    max_chars = int(w / char_width_estimate)
                    if pdf.get_string_width(text) > w and max_chars > 3:
                        text = text[:max_chars - 3] + '...'
                    elif pdf.get_string_width(text) > w:
                         text = text[:max_chars]
                else:
                    if pdf.get_string_width(text) > w:
                        text = text[:10] + '...'

                pdf.cell(w, 10, txt=text, border=1)
            pdf.ln()
        pdf.set_font("Arial", size=12)
        pdf.ln(5)
    else:
        pdf.multi_cell(page_width, 5, txt=str(content))
        pdf.ln(5)

    if fig is not None:
        pdf.add_page()
        pdf.cell(page_width, 10, txt="Visualization:", ln=True, align='L')
        pdf.ln(5)
        img_stream = io.BytesIO()
        pio.write_image(fig, img_stream, format='png', scale=2)
        img_stream.seek(0)
        temp_img_path = f"temp_plot_{uuid.uuid4()}.png"
        try:
            with open(temp_img_path, 'wb') as f:
                f.write(img_stream.getvalue())
            img_width_mm = page_width
            pdf.image(temp_img_path, x=pdf.l_margin, y=pdf.get_y(), w=img_width_mm)
        finally:
            if os.path.exists(temp_img_path):
                os.remove(temp_img_path)

    return pdf.output(dest='S').encode('latin-1')

def plot_data(df, message_key):
    fig = None
    plot_col, config_col = st.columns([3, 1])

    with config_col:
        st.text("📊 Plot Configuration")
        plot_type = st.selectbox("Plot Type", ["Bar", "Line", "Scatter", "Box", "Histogram", "Pie"], key=f"plot_type_{message_key}")
        all_columns = df.columns.tolist()
        numeric_columns = df.select_dtypes(include='number').columns.tolist()
        categorical_columns = df.select_dtypes(exclude='number').columns.tolist()

        x_axis, y_axis, values, names = None, None, None, None

        if plot_type == "Pie":
            col1, col2 = st.columns(2)
            with col1:
                values_options = numeric_columns if numeric_columns else all_columns
                if not values_options: st.warning("Need numeric data for pie values."); return None
                values = st.selectbox("Values", values_options, key=f"pie_val_{message_key}", help="Column for pie slice sizes (numeric)")
            with col2:
                names_options = categorical_columns if categorical_columns else all_columns
                if not names_options: st.warning("Need categorical data for pie labels."); return None
                names = st.selectbox("Labels", names_options, key=f"pie_name_{message_key}", help="Column for pie slice labels (categorical)")
            if not values or not names: return None

        elif plot_type == "Histogram":
            options = numeric_columns if numeric_columns else all_columns
            if not options: st.warning("Histogram requires a column."); return None
            x_axis = st.selectbox("Column", options, key=f"hist_x_{message_key}", help="Column to generate histogram from")

        else:
            col1, col2 = st.columns(2)
            with col1:
                x_options = all_columns
                if not x_options: st.warning("Need X-axis data."); return None
                x_axis = st.selectbox("X-axis", x_options, key=f"xy_x_{message_key}")
            with col2:
                y_options = numeric_columns if numeric_columns else all_columns
                if not y_options: st.warning("Need Y-axis data."); return None
                y_axis = st.selectbox("Y-axis", y_options, key=f"xy_y_{message_key}", help="Typically a numeric column for the Y-axis")
            if not x_axis or not y_axis: return None

        st.markdown("---")

        st.write("Customize Axes & Title")
        title = st.text_input("Plot Title", key=f"plot_title_{message_key}", placeholder=f"Auto: {plot_type} Plot")
        col1, col2 = st.columns(2)
        with col1:
            x_title = st.text_input("X-axis Title", key=f"plot_xtitle_{message_key}", placeholder=f"Auto: {x_axis if x_axis else 'X'}")
        with col2:
            y_title = st.text_input("Y-axis Title", key=f"plot_ytitle_{message_key}", placeholder=f"Auto: {y_axis if y_axis else 'Y'}")

    st.markdown("---")

    st.text("⬇️ Export Data & Plot")

    export_filename_base = f"chat_export_{message_key}_{time.strftime('%Y%m%d_%H%M%S')}"

    excel_buffer = io.BytesIO()
    df.to_excel(excel_buffer, index=False)
    excel_bytes = excel_buffer.getvalue()

    try:
        pptx_bytes = export_to_ppt(df, fig)
    except Exception as e:
        st.warning(f"Could not generate PowerPoint preview: {e}")
        pptx_bytes = None

    try:
        pdf_bytes = export_to_pdf(df, fig)
    except Exception as e:
        st.warning(f"Could not generate PDF preview: {e}")
        pdf_bytes = None

    col1, col2, col3 = st.columns(3)
    with col1:
        st.download_button(
            label=".xlsx",
            data=excel_bytes,
            file_name=f"{export_filename_base}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key=f"dwnld_excel_{message_key}",
            help="Download data as Excel",
            use_container_width=True
        )
    with col2:
        st.download_button(
            label=".pptx",
            data=pptx_bytes if pptx_bytes else b"",
            file_name=f"{export_filename_base}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"dwnld_ppt_{message_key}",
            help="Download data and plot as PowerPoint",
            use_container_width=True,
            disabled=pptx_bytes is None
        )
    with col3:
        st.download_button(
            label=".pdf",
            data=pdf_bytes if pdf_bytes else b"",
            file_name=f"{export_filename_base}.pdf",
            mime="application/pdf",
            key=f"dwnld_pdf_{message_key}",
            help="Download data and plot as PDF",
            use_container_width=True,
            disabled=pdf_bytes is None
        )

    with plot_col:
        try:
            plot_generated = False
            if plot_type == "Pie" and values and names:
                fig = px.pie(df, values=values, names=names)
                plot_generated = True
            elif plot_type == "Histogram" and x_axis:
                fig = px.histogram(df, x=x_axis)
                plot_generated = True
            elif x_axis and y_axis:
                plot_func = {
                    "Bar": px.bar, "Line": px.line, "Scatter": px.scatter, "Box": px.box
                }.get(plot_type)
                if plot_func:
                    fig = plot_func(df, x=x_axis, y=y_axis)
                    plot_generated = True

            if fig:
                final_title = title if title else f"{plot_type} Plot of {y_axis if y_axis else values} vs {x_axis if x_axis else names}"
                final_x_title = x_title if x_title else (x_axis if x_axis else names)
                final_y_title = y_title if y_title else (y_axis if y_axis else values)

                fig.update_layout(
                    title=final_title,
                    xaxis_title=final_x_title,
                    yaxis_title=final_y_title,
                    legend_title="Legend"
                )
                st.plotly_chart(fig, use_container_width=True)
            elif plot_generated is False and plot_type in ["Pie", "Histogram", "Bar", "Line", "Scatter", "Box"]:
                 st.info("Configure plot options on the right to generate visualization.")

        except Exception as e:
            st.error(f"Could not generate plot: {e}")
            fig = None

    return fig

def show_ingestion_results(results: List[Dict]):
    """Display ingestion results with proper formatting"""
    success_count = sum(1 for r in results if r['success'])
    failure_count = len(results) - success_count
    
    if success_count > 0:
        st.success(f"✅ Successfully ingested {success_count} file(s)")
    
    if failure_count > 0:
        st.error(f"❌ Failed to ingest {failure_count} file(s)")
    
    with st.expander("See detailed results", expanded=True):
        for result in results:
            if result['success']:
                with st.container():
                    st.markdown(f"✅ **{result['filename']}**")
                    st.markdown(f"- {result['message']}")
                    if 'metadata' in result:
                        content_type = result['metadata'].get('Content-Type', 'Unknown type')
                        st.markdown(f"- Type: `{content_type}`")
                    st.markdown("---")
            else:
                with st.container():
                    st.markdown(f"❌ **{result['filename']}**")
                    st.markdown(f"- Error: {result['error']}")
                    st.markdown("---")

def launch_bot():

    if 'cfg' not in st.session_state:
        try:
            corpus_keys = str(os.environ['corpus_keys']).split(',')
            cfg = OmegaConf.create({
                'corpus_keys': corpus_keys,
                'api_key': str(os.environ['api_key']),
                'title': os.environ.get('title', 'Chatbot'),
                'source_data_desc': os.environ.get('source_data_desc', 'AI Assistant'),
                'streaming': isTrue(os.environ.get('streaming', True)),
                'prompt_name': os.environ.get('prompt_name', None),
                'examples': os.environ.get('examples', ''),
                'language': 'English'
            })
            st.session_state.cfg = cfg
            st.session_state.vq = VectaraQuery(cfg.api_key, cfg.corpus_keys, cfg.prompt_name)
            st.session_state.device_id = str(uuid.uuid4())
            st.session_state.feedback_key = 0
            st.session_state.message_counter = 0

            st.session_state.messages = [{"role": "assistant", "content": INITIAL_ASSISTANT_MESSAGE, "avatar": ASSISTANT_AVATAR, "key": f"msg_{st.session_state.message_counter}"}]
            st.session_state.message_counter += 1

            example_messages = [example.strip() for example in cfg.examples.split(",") if example.strip()]
            st.session_state.example_messages = example_messages[:MAX_EXAMPLES]
            st.session_state.show_examples = bool(st.session_state.example_messages)
            st.session_state.prompt_input = ""
            st.session_state.pending_prompt = None
            st.session_state.selected_format = "Summarized Text"

        except KeyError as e:
            st.error(f"Missing required environment variable: {e}. Please check your .env file.")
            st.stop()
        except Exception as e:
            st.error(f"Error initializing configuration: {e}")
            st.stop()

    cfg = st.session_state.cfg
    vq = st.session_state.vq

    st.set_page_config(page_title=cfg.title, layout="wide", initial_sidebar_state="expanded")

    def reset_chat():
        st.session_state.message_counter = 0
        st.session_state.messages = [{"role": "assistant", "content": INITIAL_ASSISTANT_MESSAGE, "avatar": ASSISTANT_AVATAR, "key": f"msg_{st.session_state.message_counter}"}]
        st.session_state.message_counter += 1
        st.session_state.vq.conv_id = None
        st.session_state.show_examples = bool(st.session_state.example_messages)
        st.session_state.feedback_key += 1
        st.session_state.prompt_input = ""
        st.session_state.pending_prompt = None
        st.session_state.selected_format = "Summarized Text"

    def generate_response(question, lang_code, temperature, again=False):
        return vq.submit_query(question, lang_code, temperature, again)

    def generate_streaming_response(question, lang_code, temperature):
        return vq.submit_query_streaming(question, lang_code, temperature)

    with st.sidebar:
        if os.path.exists(LOGO_PATH):
            image = Image.open(LOGO_PATH)
            st.image(image, use_container_width=True)
        else:
            st.warning(f"Logo file not found at {LOGO_PATH}")

        st.caption(cfg.title)

        st.markdown("---")

        uploaded_files = st.file_uploader(
            label="Upload Files",
            help="Supported formats: PDF, PPT/PPTX, DOC/DOCX, Markdown, HTML",
            accept_multiple_files=True,
            type=['pdf', 'ppt', 'pptx', 'doc', 'docx', 'md', 'html', 'htm']
        )

        if uploaded_files:
            st.info(f"📁 {len(uploaded_files)} file(s) selected")
            if st.button("Ingest", key="ingest_button", use_container_width=True):
                with st.spinner("Ingesting files..."):
                    try:
                        results = vq.ingest_files_to_corpus(uploaded_files, corpus_key="custom_uploaded_data")
                        show_ingestion_results(results)
                    except Exception as e:
                        st.error(f"Failed to process files: {str(e)}")

        st.markdown("---")

        selected_language = st.selectbox(
            'Language',
            options=list(LANGUAGES.keys()),
            index=list(LANGUAGES.keys()).index(st.session_state.get('language', 'English')),
            key='language_select'
        )
        if st.session_state.get('language') != selected_language:
            st.session_state.language = selected_language
            cfg.language = selected_language
            reset_chat()
            st.rerun()
        
        temperature = st.slider("Temperature", min_value=0.0, max_value=1.0, value=0.1, step=0.1, help="How creative do you want my answers to be?")
        st.session_state.temperature = temperature

        st.markdown("---")
        if st.button('Start Over', key='start_over_button', use_container_width=True):
            reset_chat()
            st.rerun()

    st.markdown(f"<h2 style='text-align: center;'>{cfg.title}</h2>", unsafe_allow_html=True)

    chat_container = st.container()
    with chat_container:
        for i, message in enumerate(st.session_state.messages):
            with st.chat_message(message["role"], avatar=message.get("avatar")):
                message_key = message.get("key", f"msg_{i}")
                content = message["content"]

                if isinstance(content, str):
                    st.markdown(content)
                elif isinstance(content, pd.DataFrame) or (isinstance(content, dict) and 'data' in content):
                    df = content['data'] if isinstance(content, dict) else content
                    cutoff_type = content.get('cutoff_type') if isinstance(content, dict) else None
                    st.dataframe(df, use_container_width=True)
                    if not df.empty:
                        with st.expander("📊 View & Export Data Visualization", expanded=False):
                            fig = plot_data(df, message_key)

                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("Analyze & Provide Insights", key=f"analyze_{message_key}", use_container_width=True):
                                original_query = next(
                                    (msg["content"] for msg in reversed(st.session_state.messages[:i])
                                    if msg["role"] == "user"),
                                    "Unknown Query"
                                )

                                analysis = vq.analyze_data_with_claude(df, original_query)

                                st.session_state.messages.append({
                                    "role": "assistant",
                                    "content": analysis,
                                    "avatar": ASSISTANT_AVATAR,
                                    "key": f"msg_{st.session_state.message_counter}"
                                })
                                st.session_state.message_counter += 1
                                st.rerun()
                        with col2:
                            full_result_button = st.button(
                                "Provide Full Result",
                                key=f"full_result_{message_key}",
                                disabled=cutoff_type == 'no_significant_drop',
                                use_container_width=True
                            )
                            if full_result_button:
                                original_query = next(
                                    (msg["content"] for msg in reversed(st.session_state.messages[:i])
                                    if msg["role"] == "user"),
                                    "Unknown Query"
                                )

                                if "[table]" not in original_query.lower():
                                    original_query = f"[table] {original_query}"

                                st.session_state.pending_prompt = original_query
                                st.session_state.selected_format_for_prompt = "Data Tables & Visualisation"
                                st.session_state.show_full_result = True
                                st.rerun()

            
                else:
                    st.write(content)

                if (i == len(st.session_state.messages) - 1 and
                    message["role"] == "assistant" and
                    isinstance(content, str) and
                    content != INITIAL_ASSISTANT_MESSAGE):

                    with st.popover("⬇️ Export Text", use_container_width=False):
                        export_filename_base = f"chat_export_text_{message_key}_{time.strftime('%Y%m%d_%H%M%S')}"
                        try:
                            pptx_text_bytes = export_to_ppt(content)
                        except Exception as e:
                             st.warning(f"PPT export failed: {e}")
                             pptx_text_bytes = None
                        try:
                            pdf_text_bytes = export_to_pdf(content)
                        except Exception as e:
                             st.warning(f"PDF export failed: {e}")
                             pdf_text_bytes = None

                        col1, col2 = st.columns(2)
                        with col1:
                            st.download_button(
                                label=".pptx", data=pptx_text_bytes if pptx_text_bytes else b"",
                                file_name=f"{export_filename_base}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"dwnld_text_ppt_{message_key}",
                                help="Download text as PowerPoint",
                                use_container_width=True,
                                disabled=pptx_text_bytes is None
                            )
                        with col2:
                            st.download_button(
                                label=".pdf", data=pdf_text_bytes if pdf_text_bytes else b"",
                                file_name=f"{export_filename_base}.pdf",
                                mime="application/pdf",
                                key=f"dwnld_text_pdf_{message_key}",
                                help="Download text as PDF",
                                use_container_width=True,
                                disabled=pdf_text_bytes is None
                            )

    if st.session_state.get('pending_prompt'):
        prompt_to_process = st.session_state.pending_prompt
        selected_format_for_pending = st.session_state.get('selected_format_for_prompt', 'Summarized Text')

        st.session_state.pending_prompt = None
        st.session_state.selected_format_for_prompt = None

        user_message_content = prompt_to_process
        if selected_format_for_pending == "Data Tables & Visualisation" and "[table]" not in prompt_to_process.lower():
            user_message_content = f"[table] {prompt_to_process}"

        lang_code = LANGUAGES.get(st.session_state.language, 'eng')

        response_generation_area = st.container()
        message_content = None

        try:
            if '[table]' in user_message_content.lower():
                 with response_generation_area:
                    with st.status("Generating Data Tables...", expanded=True):
                        st.write("Querying Data Store...")
                        again = st.session_state.get('show_full_result', False)
                        response = generate_response(user_message_content, lang_code, temperature, again)
                        st.write("Formatting Data Tables...")

                        if 'show_full_result' in st.session_state:
                            del st.session_state.show_full_result

                        if not isinstance(response, pd.DataFrame):
                             try:
                                  original_response_for_fallback = response
                                  if isinstance(response, str):
                                       parsed = False
                                       for sep in [',', ';', '\t']:
                                           try:
                                               response_df = pd.read_csv(io.StringIO(response), sep=sep, comment=None, skipinitialspace=True)
                                               if response_df.shape[1] > 0:
                                                    response = response_df
                                                    parsed = True
                                                    break
                                           except (pd.errors.ParserError, ValueError):
                                               continue
                                       if not parsed:
                                           st.warning("Could not parse string response as CSV with common delimiters. Displaying raw string.")
                                           response = original_response_for_fallback

                                  elif isinstance(response, list) and response:
                                       if isinstance(response[0], dict):
                                           response = pd.DataFrame(response)
                                       elif isinstance(response[0], list):
                                            is_header = False
                                            if len(response) > 1:
                                                first_row_is_str = all(isinstance(x, str) for x in response[0])
                                                second_row_mixed_types = not all(isinstance(x, type(response[1][0])) for x in response[1]) if len(response[1]) > 0 else False
                                                if first_row_is_str and (len(response[0]) == len(response[1]) or second_row_mixed_types):
                                                    is_header = True

                                            if is_header:
                                                response = pd.DataFrame(response[1:], columns=response[0])
                                            else:
                                                response = pd.DataFrame(response)
                                       else:
                                            st.warning("List format not recognized (expected list of dicts or list of lists). Displaying raw list.")
                                            response = original_response_for_fallback

                                  else:
                                       st.warning("Received non-tabular data for table request. Displaying raw response.")
                                       response = original_response_for_fallback

                             except Exception as parse_error:
                                  st.warning(f"Could not parse response into a table: {parse_error}. Displaying raw response.")
                                  response = original_response_for_fallback

                        message_content = response

            elif cfg.streaming:
                 with response_generation_area:
                     with st.chat_message("assistant", avatar=ASSISTANT_AVATAR):
                         stream_placeholder = st.empty()
                         stream = generate_streaming_response(user_message_content, lang_code, temperature)
                         response_str = stream_placeholder.write_stream(stream)
                 message_content = escape_dollars_outside_latex(response_str)

            else:
                with response_generation_area:
                     with st.status("Thinking...", expanded=True):
                         response_str = generate_response(user_message_content, lang_code, temperature)
                message_content = escape_dollars_outside_latex(response_str)

            if message_content is not None:
                st.session_state.messages.append({"role": "assistant", "content": message_content, "avatar": ASSISTANT_AVATAR, "key": f"msg_{st.session_state.message_counter}"})
                st.session_state.message_counter += 1

                send_amplitude_data(
                    user_query=user_message_content,
                    chat_response=str(message_content),
                    demo_name=cfg["title"],
                    language=st.session_state.language
                )
                st.session_state.feedback_key += 1
            else:
                 st.error("Failed to generate a valid response.")

            response_generation_area.empty()
            st.rerun()

        except Exception as e:
            st.error(f"An error occurred during response generation: {e}")
            response_generation_area.empty()
            error_message = f"Sorry, an error occurred: {e}"
            st.session_state.messages.append({"role": "assistant", "content": error_message, "avatar": ASSISTANT_AVATAR, "key": f"msg_{st.session_state.message_counter}"})
            st.session_state.message_counter += 1
            st.session_state.feedback_key += 1
            st.rerun()

    if st.session_state.get('show_examples', False):
        selected_example = pills("✨ Queries to Try:", st.session_state.example_messages, index=None, key="example_pills")
        if selected_example:
            st.session_state.show_examples = False
            st.session_state.pending_prompt = selected_example
            st.session_state.selected_format_for_prompt = st.session_state.get('selected_format', 'Summarized Text')
            user_message_content = selected_example
            if st.session_state.selected_format_for_prompt == "Data Tables & Visualisation" and "[table]" not in selected_example.lower():
                 user_message_content = f"[table] {selected_example}"
            st.session_state.messages.append({"role": "user", "content": user_message_content, "avatar": USER_AVATAR, "key": f"msg_{st.session_state.message_counter}"})
            st.session_state.message_counter += 1
            st.rerun()

    current_format_index = ["Summarized Text", "Data Tables & Visualisation"].index(st.session_state.get('selected_format', 'Summarized Text'))
    selected_format = pills("Response format:", ["Summarized Text", "Data Tables & Visualisation"], index=current_format_index, key="response_format_pills_main")

    if selected_format != st.session_state.get('selected_format', 'Summarized Text'):
         st.session_state.selected_format = selected_format

    if prompt := st.chat_input(f"Enter the query for {st.session_state.selected_format}", key="chat_input_main"):
        st.session_state.show_examples = False

        st.session_state.pending_prompt = prompt
        st.session_state.selected_format_for_prompt = st.session_state.selected_format

        user_message_content = prompt
        if st.session_state.selected_format == "Data Tables & Visualisation" and "[table]" not in prompt.lower():
             user_message_content = f"[table] {prompt}"
        st.session_state.messages.append({"role": "user", "content": user_message_content, "avatar": USER_AVATAR, "key": f"msg_{st.session_state.message_counter}"})
        st.session_state.message_counter += 1

        st.rerun()

if __name__ == "__main__":
    launch_bot()
